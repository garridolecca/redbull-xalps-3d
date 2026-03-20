"""
Generate a professional Esri + Red Bull X-Alps PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
RB_RED = RGBColor(0xDB, 0x0A, 0x40)
RB_DARK = RGBColor(0x0A, 0x0A, 0x18)
RB_NAVY = RGBColor(0x14, 0x14, 0x24)
ESRI_BLUE = RGBColor(0x00, 0x79, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xBB)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
CYAN = RGBColor(0x48, 0xB0, 0xC7)
GREEN = RGBColor(0x34, 0xC7, 0x59)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=RB_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_shape_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_slide(slide, left, top, width, height, items, size=16, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)  # top accent
add_shape_rect(slide, 0, 7.42, 13.333, 0.08, ESRI_BLUE)  # bottom accent

add_text(slide, 1, 1.8, 11, 1.2, "RED BULL X-ALPS 2025", size=48, color=RB_RED, bold=True)
add_text(slide, 1, 3.0, 11, 0.8, "3D Race Replay & Thermal Analysis", size=28, color=WHITE)
add_text(slide, 1, 4.0, 11, 0.6, "Interactive Geospatial Dashboard", size=20, color=LIGHT_GRAY)

add_text(slide, 1, 5.5, 5, 0.5, "Powered by ArcGIS Maps SDK for JavaScript", size=14, color=ESRI_BLUE)
add_text(slide, 1, 5.9, 5, 0.5, "Calcite Design System  |  3D SceneView", size=14, color=LIGHT_GRAY)

add_text(slide, 8, 5.5, 4, 0.5, "Data: Naviter GPS Telemetry", size=14, color=GOLD, align=PP_ALIGN.RIGHT)
add_text(slide, 8, 5.9, 4, 0.5, "37.5M GPS fixes  |  35 Athletes  |  17 Days", size=14, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════
# SLIDE 2: THE RACE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "THE RACE", size=36, color=RB_RED, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.5, "What is Red Bull X-Alps?", size=20, color=WHITE, bold=True)

add_bullet_slide(slide, 0.8, 1.8, 5.5, 5, [
    "Adventure race across the European Alps",
    "Athletes travel from Salzburg, Austria to Monaco",
    "Only two modes: paragliding or hiking on foot",
    "Athletes carry their own paraglider (~30 kg pack)",
    "No motorized transport allowed",
    "First to reach the finish line wins",
    "Race spans approximately 12-14 days",
    "Mandatory rest periods enforced at night",
    "35 athletes from 20+ countries compete",
], size=15)

add_text(slide, 7, 1.8, 5.5, 0.5, "Key Stats from 2025 Data", size=18, color=GOLD, bold=True)
# Stats cards
for i, (val, lbl) in enumerate([
    ("37.5M", "GPS Fixes Recorded"),
    ("35", "Athletes Tracked"),
    ("17", "Race Days"),
    ("4", "Sensor Types"),
    ("18", "Fields Per Fix"),
]):
    y = 2.5 + i * 0.85
    add_shape_rect(slide, 7.2, y, 5, 0.7, RB_NAVY)
    add_text(slide, 7.4, y + 0.05, 1.5, 0.6, val, size=24, color=RB_RED, bold=True)
    add_text(slide, 9, y + 0.12, 3, 0.5, lbl, size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════
# SLIDE 3: DATA PIPELINE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "DATA PIPELINE", size=36, color=RB_RED, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.5, "From GPS Instruments to 3D Visualization", size=20, color=WHITE, bold=True)

# Pipeline steps
steps = [
    ("Naviter Omni\nGPS Device", "On each athlete\n18 fields/fix, 1Hz", RB_RED),
    ("Google\nPub/Sub", "Real-time message\ningestion", RGBColor(0xFF, 0x8C, 0x00)),
    ("Google\nBigQuery", "Raw storage\n14 daily CSV exports", RGBColor(0xFF, 0xA5, 0x00)),
    ("Python\nETL Script", "Parse, deduplicate\n90s time buckets", GOLD),
    ("JSON\n17 MB", "294K tracks\n59K climb events", GREEN),
    ("ArcGIS\nSceneView", "3D terrain\nCalcite Design", ESRI_BLUE),
]
for i, (title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.1
    add_shape_rect(slide, x, 2.3, 1.8, 1.4, RB_NAVY)
    # colored top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.3), Inches(1.8), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, x + 0.1, 2.5, 1.6, 0.6, title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 3.1, 1.6, 0.5, desc, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(slide, x + 1.8, 2.7, 0.3, 0.5, ">", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Processing details
add_text(slide, 0.8, 4.2, 11, 0.5, "Processing Details", size=18, color=WHITE, bold=True)
add_bullet_slide(slide, 0.8, 4.8, 5.5, 2.5, [
    "Sensor priority: Naviter > Satellite > Flarm > OGN",
    "Time-bucket deduplication: 90-second intervals",
    "Coordinate precision: 4 decimal places (~11m)",
    "Altitude rounded to integers",
    "Timestamps filtered to race window only",
], size=13)
add_bullet_slide(slide, 7, 4.8, 5.5, 2.5, [
    "Input: 14 CSV files, 1.16M rows, 37.5M data points",
    "Output: 294K tracks + 59K climb + 62K wind",
    "File size reduced: 25 MB > 17 MB",
    "Streaming download with progress indicator",
    "JSON parsed client-side in browser",
], size=13)

# ═══════════════════════════════════════════════
# SLIDE 4: RAW DATA SCHEMA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "RAW DATA SCHEMA", size=36, color=RB_RED, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.5, "18 Fields Per GPS Fix (from Naviter Omni Data Dictionary)", size=20, color=WHITE, bold=True)

fields = [
    ("GPS Timestamp", "Unix UTC", "Position in time"),
    ("Latitude / Longitude", "WGS 84", "Position on Earth"),
    ("GPS Altitude", "Meters (ellipsoid)", "Height from GPS"),
    ("Baro Altitude", "Meters (QNE)", "Height from barometer"),
    ("Altitude Above Ground", "Meters (SRTM)", "Height above terrain"),
    ("Speed / Groundspeed", "km/h", "Movement rate"),
    ("Heading", "Degrees", "Direction of travel"),
    ("Vertical Speed", "m/s", "Rate of climb or descent"),
    ("Windspeed", "km/h", "Measured by instrument"),
    ("Wind Direction", "Degrees", "Where wind comes from"),
    ("UserActivity", "S / F / H", "Stationary / Flying / Hiking"),
    ("Airspace", "String", "Airspace zone identifier"),
    ("Rest Status", "0 / 1", "Mandatory rest active"),
    ("Nightpass Status", "0 / 1", "Night movement allowed"),
    ("GPS Quality", "0-50 SNR", "Signal quality indicator"),
]
for i, (name, unit, desc) in enumerate(fields):
    y = 1.8 + i * 0.35
    col = i % 2
    x = 0.8 + col * 6.2
    if i >= 8:
        y = 1.8 + (i - 8) * 0.35
        x = 7
    else:
        x = 0.8
    bg_color = RB_NAVY if (i // 1) % 2 == 0 else RB_DARK
    add_text(slide, x, y, 2.2, 0.3, name, size=11, color=WHITE, bold=True)
    add_text(slide, x + 2.3, y, 1.2, 0.3, unit, size=10, color=CYAN)
    add_text(slide, x + 3.6, y, 2.2, 0.3, desc, size=10, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════
# SLIDE 5: ANALYSIS METHODOLOGY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "ANALYSIS METHODOLOGY", size=36, color=RB_RED, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.5, "How Indicators Are Derived from Raw Data", size=20, color=WHITE, bold=True)

# Climb events
add_shape_rect(slide, 0.6, 1.9, 5.8, 2.4, RB_NAVY)
add_text(slide, 0.8, 2.0, 5, 0.4, 'Climb Events (labeled "Thermals")', size=16, color=GOLD, bold=True)
add_bullet_slide(slide, 0.8, 2.5, 5.4, 1.8, [
    "Filter: Vertical Speed > 1.0 m/s",
    "Source: Naviter sensor only (type = 'N')",
    "Sampling: 90-second time buckets per athlete",
    "No clustering into distinct thermal events",
    "UserActivity field NOT checked (may include hiking)",
    "Threshold of 1.0 m/s is arbitrary, not a standard",
], size=12)

# Wind
add_shape_rect(slide, 6.8, 1.9, 5.8, 2.4, RB_NAVY)
add_text(slide, 7.0, 2.0, 5, 0.4, "Wind Rose Analysis", size=16, color=CYAN, bold=True)
add_bullet_slide(slide, 7.0, 2.5, 5.4, 1.8, [
    "Source: Windspeed + Wind Direction from Naviter Omni",
    "These are onboard instrument readings, not weather data",
    "Sampling: 5-minute intervals per athlete",
    "16 directional bins, 4 speed bins (0-5, 5-15, 15-25, 25+)",
    "Time window: +/- 1.5 hours of current playback time",
    "Stats: average speed, max speed, dominant direction",
], size=12)

# Altitude
add_shape_rect(slide, 0.6, 4.6, 5.8, 1.8, RB_NAVY)
add_text(slide, 0.8, 4.7, 5, 0.4, "Altitude Profile", size=16, color=GREEN, bold=True)
add_bullet_slide(slide, 0.8, 5.2, 5.4, 1.2, [
    "Plots raw GPS Altitude (ellipsoid height, meters)",
    "Full track timeline for solo'd athlete",
    "Red marker shows current playback position",
    "Canvas-drawn sparkline, updates during playback",
], size=12)

# Caveats
add_shape_rect(slide, 6.8, 4.6, 5.8, 1.8, RB_NAVY)
add_text(slide, 7.0, 4.7, 5, 0.4, "Important Caveats", size=16, color=RB_RED, bold=True)
add_bullet_slide(slide, 7.0, 5.2, 5.4, 1.2, [
    "\"Thermals\" = inferred from vertical speed, not measured",
    "Positive VS could be ridge lift, wave, or GPS noise",
    "Wind data only from Naviter, not all sensors",
    "Athlete names approximate (bib mapping may vary)",
], size=12)

# ═══════════════════════════════════════════════
# SLIDE 6: VISUALIZATION FEATURES
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "VISUALIZATION FEATURES", size=36, color=RB_RED, bold=True)

features = [
    ("3D Race Replay", "Animated athlete markers moving along GPS tracks over Alpine terrain with real-time interpolation", RB_RED),
    ("Numbered Markers", "Each athlete shows their bib number in a colored circle with callout line to terrain", RGBColor(0xFF, 0x8C, 0x00)),
    ("Track Lines", "Dual-layer glow effect: translucent outer + bright inner core at actual GPS altitude", GOLD),
    ("Wind Rose Chart", "Canvas-drawn polar diagram with 16 directions, 4 speed bins, live stats", CYAN),
    ("Thermal Statistics", "Active count, average and max climb rate, updating every 500ms during playback", GREEN),
    ("Altitude Profile", "Sparkline elevation chart for solo'd athlete with moving time marker", GREEN),
    ("Solo Mode", "Filter all data to a single athlete: tracks, climb events, wind, altitude profile", RGBColor(0xAF, 0x52, 0xDE)),
    ("Dynamic Lighting", "Sun position updates with playback time for realistic day/night transitions", RGBColor(0xFF, 0xD7, 0x00)),
]
for i, (title, desc, color) in enumerate(features):
    row = i // 2
    col = i % 2
    x = 0.6 + col * 6.2
    y = 1.4 + row * 1.45
    add_shape_rect(slide, x, y, 5.8, 1.2, RB_NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, x + 0.2, y + 0.1, 5.2, 0.4, title, size=15, color=WHITE, bold=True)
    add_text(slide, x + 0.2, y + 0.55, 5.2, 0.6, desc, size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════
# SLIDE 7: TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_text(slide, 0.8, 0.4, 11, 0.8, "TECHNICAL ARCHITECTURE", size=36, color=RB_RED, bold=True)

# Left column
add_text(slide, 0.8, 1.4, 5, 0.4, "Frontend Stack", size=18, color=ESRI_BLUE, bold=True)
add_bullet_slide(slide, 0.8, 1.9, 5.5, 3, [
    "ArcGIS Maps SDK for JavaScript 4.31",
    "Calcite Design System 2.13.2 (dark mode)",
    "SceneView with world elevation service",
    "GraphicsLayer with absolute-height elevation",
    "LineSymbol3D with dual-layer glow effect",
    "IconSymbol3DLayer with inline SVG markers",
    "PointSymbol3D with callout lines",
    "HTML5 Canvas for wind rose + altitude charts",
    "Streaming fetch with download progress",
], size=13)

# Right column
add_text(slide, 7, 1.4, 5, 0.4, "Performance Optimizations", size=18, color=GOLD, bold=True)
add_bullet_slide(slide, 7, 1.9, 5.5, 3, [
    "Binary search for time-windowed data queries",
    "Throttled analysis updates (every 500ms)",
    "requestIdleCallback for non-critical DOM updates",
    "No map graphics for thermals/wind (charts only)",
    "Animation: only positions + time display per frame",
    "90-second time buckets reduce data 3x",
    "4 decimal coordinate precision saves ~30% JSON size",
    "Auto-skip empty time gaps during playback",
    "Shadows + ambient occlusion disabled for GPU perf",
], size=13)

add_text(slide, 0.8, 5.2, 5, 0.4, "Data Processing", size=18, color=GREEN, bold=True)
add_bullet_slide(slide, 0.8, 5.7, 5.5, 1.5, [
    "Python script: csv, json, defaultdict",
    "Sensor priority deduplication (N > S > F > O)",
    "Streaming CSV reader with 10MB field limit",
    "Output: single JSON file served via GitHub Pages",
], size=13)

add_text(slide, 7, 5.2, 5, 0.4, "Deployment", size=18, color=RB_RED, bold=True)
add_bullet_slide(slide, 7, 5.7, 5.5, 1.5, [
    "Single HTML file + JSON data file",
    "Hosted on GitHub Pages (auto-deploy on push)",
    "No server required - fully client-side",
    "Login gate + intro guide for onboarding",
], size=13)

# ═══════════════════════════════════════════════
# SLIDE 8: THANK YOU
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, RB_RED)
add_shape_rect(slide, 0, 7.42, 13.333, 0.08, ESRI_BLUE)

add_text(slide, 1, 2.2, 11, 1, "THANK YOU", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.4, 11, 0.6, "Red Bull X-Alps 2025 | 3D Race Replay & Thermal Analysis", size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 1, 4.8, 11, 0.5, "github.com/garridolecca/redbull-xalps-3d", size=16, color=ESRI_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.3, 11, 0.5, "garridolecca.github.io/redbull-xalps-3d", size=16, color=RB_RED, align=PP_ALIGN.CENTER)

add_text(slide, 1, 6.2, 5, 0.4, "Powered by Esri ArcGIS", size=14, color=ESRI_BLUE)
add_text(slide, 7, 6.2, 5, 0.4, "Data by Red Bull X-Alps / Naviter", size=14, color=GOLD, align=PP_ALIGN.RIGHT)

# Save
out = r"E:\3. Data\Red Bull\app\RedBull_XAlps_2025_Dashboard.pptx"
prs.save(out)
print(f"Saved: {out}")
